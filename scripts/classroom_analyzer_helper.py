#!/usr/bin/env python3
import argparse
import sys
import os
import json
import subprocess
from pathlib import Path

def fetch_audio(url: str, output_path: Path) -> int:
    """Downloads audio from a YouTube/video URL and converts to MP3 using yt-dlp."""
    print(f"Downloading audio from: {url}")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    # yt-dlp command to extract audio as mp3
    cmd = [
        "yt-dlp",
        "-x",
        "--audio-format", "mp3",
        "--audio-quality", "0",
        "-o", str(output_path.with_suffix("")),  # yt-dlp appends .mp3 automatically
        url
    ]
    
    try:
        result = subprocess.run(cmd, check=True, capture_output=True, text=True)
        # Check if the output file exists (yt-dlp sometimes names it slightly differently if extension is appended)
        expected_file = output_path.with_suffix(".mp3")
        if expected_file.exists() and expected_file != output_path:
            if output_path.exists():
                os.remove(output_path)
            os.rename(expected_file, output_path)
        
        if not output_path.exists():
            # Search for any mp3 in same dir matching the name prefix
            prefix = output_path.stem
            for f in output_path.parent.glob(f"{prefix}*.mp3"):
                os.rename(f, output_path)
                break
                
        if output_path.exists():
            print(f"[OK] Downloaded and converted audio to: {output_path}")
            return 0
        else:
            print(f"Error: Output file {output_path} was not created.", file=sys.stderr)
            print(f"yt-dlp stdout: {result.stdout}", file=sys.stderr)
            print(f"yt-dlp stderr: {result.stderr}", file=sys.stderr)
            return 1
    except subprocess.CalledProcessError as e:
        print(f"Error executing yt-dlp: {e}", file=sys.stderr)
        print(f"yt-dlp stderr:\n{e.stderr}", file=sys.stderr)
        return 1
    except Exception as e:
        print(f"Unexpected error: {e}", file=sys.stderr)
        return 1

def transcribe_audio(audio_path: Path, output_srt: Path, output_txt: Path, model_name: str, device: str) -> int:
    """Transcribes audio using faster-whisper and outputs SRT and TXT files."""
    print(f"Loading Whisper model '{model_name}' on {device}...")
    try:
        from faster_whisper import WhisperModel
    except ImportError:
        print("Error: 'faster-whisper' package is not installed. Please run with 'uv run --with faster-whisper'.", file=sys.stderr)
        return 1

    try:
        # Load model with fallbacks if medium fails
        model = None
        for current_model in [model_name, "small", "base"]:
            try:
                model = WhisperModel(current_model, device=device, compute_type="int8")
                print(f"Successfully loaded model: {current_model}")
                break
            except Exception as ex:
                print(f"Failed to load Whisper model '{current_model}': {ex}. Trying next fallback...", file=sys.stderr)
        
        if not model:
            print("Error: Could not load any Whisper model.", file=sys.stderr)
            return 1

        print(f"Transcribing audio file: {audio_path}")
        segments, info = model.transcribe(
            str(audio_path),
            language="zh",
            beam_size=5,
            word_timestamps=False
        )

        srt_lines = []
        txt_lines = []
        
        def format_time(seconds: float) -> str:
            h = int(seconds // 3600)
            m = int((seconds % 3600) // 60)
            s = int(seconds % 60)
            ms = int((seconds - int(seconds)) * 1000)
            return f"{h:02d}:{m:02d}:{s:02d},{ms:03d}"

        print("Processing segments...")
        for i, seg in enumerate(segments, 1):
            start_str = format_time(seg.start)
            end_str = format_time(seg.end)
            text = seg.text.strip()
            
            srt_lines.append(f"{i}\n{start_str} --> {end_str}\n{text}\n")
            txt_lines.append(text)
            
        output_srt.parent.mkdir(parents=True, exist_ok=True)
        output_srt.write_text("\n".join(srt_lines), encoding="utf-8")
        
        output_txt.parent.mkdir(parents=True, exist_ok=True)
        output_txt.write_text("\n".join(txt_lines), encoding="utf-8")
        
        print(f"[OK] Subtitles written to: {output_srt}")
        print(f"[OK] Full transcript written to: {output_txt}")
        return 0
    except Exception as e:
        print(f"Error during transcription: {e}", file=sys.stderr)
        return 1

def generate_slides(analysis_path: Path, output_pptx: Path, style: str) -> int:
    """Generates a PPTX presentation based on the analysis text file."""
    print(f"Generating PPTX slides with style '{style}'...")
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        print("Error: 'python-pptx' package is not installed. Please run with 'uv run --with python-pptx'.", file=sys.stderr)
        return 1

    if not analysis_path.exists():
        print(f"Error: Analysis file {analysis_path} does not exist.", file=sys.stderr)
        return 1

    try:
        analysis_text = analysis_path.read_text(encoding="utf-8")
    except Exception as e:
        print(f"Error reading analysis file: {e}", file=sys.stderr)
        return 1

    # Simple outline parser (looks for headings and bullet points)
    slides_data = []
    current_slide = None
    
    for line in analysis_text.splitlines():
        line = line.strip()
        if not line:
            continue
        if line.startswith("# ") or line.startswith("## ") or line.startswith("### "):
            # Heading -> New Slide
            title = line.replace("#", "").strip()
            if current_slide:
                slides_data.append(current_slide)
            current_slide = {"title": title, "bullets": []}
        elif line.startswith("-") or line.startswith("*") or line.startswith("1."):
            # Bullet point
            bullet = line.lstrip("-*1. ").strip()
            if current_slide:
                current_slide["bullets"].append(bullet)
        elif current_slide and len(line) > 10:
            # Descriptive text as sub-bullet or paragraph
            current_slide["bullets"].append(line)

    if current_slide:
        slides_data.append(current_slide)

    # Ensure at least 12 slides by adding dummy slides if necessary, or breaking down slides
    expanded_slides = []
    for s in slides_data:
        if len(s["bullets"]) > 5:
            # Break down slide into two parts
            mid = len(s["bullets"]) // 2
            expanded_slides.append({"title": s["title"] + " (I)", "bullets": s["bullets"][:mid]})
            expanded_slides.append({"title": s["title"] + " (II)", "bullets": s["bullets"][mid:]})
        else:
            expanded_slides.append(s)
            
    while len(expanded_slides) < 12:
        expanded_slides.append({
            "title": f"學習共同體課例探究與省思 - 專題討論 ({len(expanded_slides) + 1})",
            "bullets": [
                "聚焦微觀對話中的聆聽關係建立",
                "探究協同合作與伸展跳躍學習的脈絡",
                "思考教師如何建立公共溝通民主語言"
            ]
        })


def generate_slides(analysis_path: Path, output_pptx: Path, style: str) -> int:
    """Generates a PPTX presentation based on the analysis text file."""
    print(f"Generating PPTX slides with style '{style}'...")
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        print("Error: 'python-pptx' package is not installed. Please run with 'uv run --with python-pptx'.", file=sys.stderr)
        return 1

    if not analysis_path.exists():
        print(f"Error: Analysis file {analysis_path} does not exist.", file=sys.stderr)
        return 1

    try:
        analysis_text = analysis_path.read_text(encoding="utf-8")
    except Exception as e:
        print(f"Error reading analysis file: {e}", file=sys.stderr)
        return 1

    # Simple outline parser (looks for headings and bullet points)
    slides_data = []
    current_slide = None
    
    for line in analysis_text.splitlines():
        line = line.strip()
        if not line:
            continue
        if line.startswith("# ") or line.startswith("## ") or line.startswith("### "):
            # Heading -> New Slide
            title = line.replace("#", "").strip().replace("**", "").replace("*", "")
            if current_slide:
                slides_data.append(current_slide)
            current_slide = {"title": title, "bullets": []}
        elif line.startswith("-") or line.startswith("*") or line.startswith("1."):
            # Bullet point
            bullet = line.lstrip("-*1. ").strip().replace("**", "").replace("*", "")
            if current_slide:
                current_slide["bullets"].append(bullet)
        elif current_slide and len(line) > 10:
            # Descriptive text as sub-bullet or paragraph
            bullet = line.strip().replace("**", "").replace("*", "")
            current_slide["bullets"].append(bullet)

    if current_slide:
        slides_data.append(current_slide)

    # Ensure at least 12 slides by adding dummy slides if necessary, or breaking down slides
    expanded_slides = []
    for s in slides_data:
        if len(s["bullets"]) > 5:
            # Break down slide into two parts
            mid = len(s["bullets"]) // 2
            expanded_slides.append({"title": s["title"] + " (I)", "bullets": s["bullets"][:mid]})
            expanded_slides.append({"title": s["title"] + " (II)", "bullets": s["bullets"][mid:]})
        else:
            expanded_slides.append(s)
            
    while len(expanded_slides) < 12:
        expanded_slides.append({
            "title": f"學習共同體課例探究與省思 - 專題討論 ({len(expanded_slides) + 1})",
            "bullets": [
                "聚焦微觀對話中的聆聽關係建立",
                "探究協同合作與伸展跳躍學習的脈絡",
                "思考教師如何建立公共溝通民主語言"
            ]
        })

    # Create Presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Styles definitions
    colors = {
        "pastel": {"bg": RGBColor(245, 245, 240), "title": RGBColor(80, 50, 40), "text": RGBColor(60, 60, 60)},
        "blue": {"bg": RGBColor(230, 240, 250), "title": RGBColor(10, 40, 80), "text": RGBColor(30, 50, 80)},
        "modern": {"bg": RGBColor(250, 250, 250), "title": RGBColor(30, 30, 30), "text": RGBColor(80, 80, 80)},
        "learning": {"bg": RGBColor(238, 242, 235), "title": RGBColor(35, 65, 45), "text": RGBColor(75, 80, 75)}
    }
    style_config = colors.get(style, colors["modern"])

    for i, slide_info in enumerate(expanded_slides[:15]):  # limit to max 15 slides
        # Use blank layout
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Apply background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = style_config["bg"]
        
        # Title Box (using 16:9 dimensions)
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.733), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_info["title"]
        p.font.name = "Microsoft JhengHei"
        
        # Determine title size based on length to fit on a single line
        title_text = slide_info["title"]
        title_len = len(title_text)
        if title_len <= 15:
            title_size = 40
        elif title_len <= 20:
            title_size = 36
        else:
            title_size = 32
            
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = style_config["title"]
        
        # Determine if we should insert an illustration on this slide
        illustration_path = None
        images_dir = Path("output/images")
        p_img = images_dir / f"slide_{i+1}.png"
        if p_img.exists():
            illustration_path = str(p_img)

        # Content Box Width: default is 11.733 inches if no illustration, or 6.5 inches if illustration
        content_width = Inches(6.5) if illustration_path else Inches(11.733)
        
        # Content Box
        contentBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), content_width, Inches(5.0))
        c_tf = contentBox.text_frame
        c_tf.word_wrap = True
        
        # Determine text size based on number of bullets and total characters to fit on page
        bullets = slide_info["bullets"][:6]
        num_bullets = len(bullets)
        total_chars = sum(len(b) for b in bullets)
        
        if num_bullets <= 3 and total_chars < 150:
            text_size = 24
        elif num_bullets <= 4 and total_chars < 220:
            text_size = 22
        else:
            text_size = 20
            
        for bullet in bullets:
            p_b = c_tf.add_paragraph()
            p_b.text = "• " + bullet
            p_b.font.name = "Microsoft JhengHei"
            p_b.font.size = Pt(text_size)
            p_b.font.color.rgb = style_config["text"]
            p_b.space_after = Pt(10)

        # Insert illustration if applicable (on the right side: X=8.3, Y=2.0, W=4.2, H=4.2)
        if illustration_path:
            try:
                slide.shapes.add_picture(illustration_path, Inches(8.3), Inches(2.0), Inches(4.2), Inches(4.2))
                print(f"[OK] Inserted illustration {illustration_path} on slide {i+1}")
            except Exception as img_err:
                print(f"Error inserting illustration on slide {i+1}: {img_err}", file=sys.stderr)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(output_pptx))
        print(f"[OK] Created PowerPoint slides: {output_pptx}")
    except PermissionError:
        output_pptx = output_pptx.parent / (output_pptx.stem + "_copy.pptx")
        print(f"[WARNING] Permission denied on saving PPTX (is slides.pptx open in PowerPoint?). Saving to copy instead: {output_pptx}", file=sys.stderr)
        prs.save(str(output_pptx))
        print(f"[OK] Created PowerPoint slides: {output_pptx}")
    
    # Automatically generate corresponding HTML interactive slides
    output_html = output_pptx.with_suffix(".html")
    try:
        generate_html_slides(expanded_slides, output_html, style)
    except PermissionError:
        output_html = output_html.parent / (output_html.stem + "_copy.html")
        print(f"[WARNING] Permission denied on saving HTML. Saving to copy instead: {output_html}", file=sys.stderr)
        generate_html_slides(expanded_slides, output_html, style)
    return 0


def generate_html_slides(expanded_slides: list, output_html: Path, style: str) -> int:
    """Generates an HTML presentation slide deck."""
    print(f"Generating HTML slides with style '{style}'...")
    
    # Styles definitions in hex
    colors_hex = {
        "pastel": {"bg": "#F5F5F0", "title": "#503228", "text": "#3C3C3C"},
        "blue": {"bg": "#E6F0FA", "title": "#0A2850", "text": "#1E3250"},
        "modern": {"bg": "#FAFAFA", "title": "#1E1E1E", "text": "#505050"},
        "learning": {"bg": "#EEF2EB", "title": "#23412D", "text": "#4B504B"}
    }
    style_config = colors_hex.get(style, colors_hex["modern"])
    
    slides_html_list = []
    
    for i, slide_info in enumerate(expanded_slides[:15]):
        title = slide_info["title"]
        bullets = slide_info["bullets"][:6]
        
        # Check if illustration exists
        illustration_src = f"./images/slide_{i+1}.png"
        images_dir = output_html.parent / "images"
        p_img = images_dir / f"slide_{i+1}.png"
        
        has_image = p_img.exists()
        
        bullets_html = "\n".join(f"          <li>{b}</li>" for b in bullets)
        
        image_html = ""
        if has_image:
            image_html = f"""
        <div class="slide-image">
          <img src="{illustration_src}" alt="Slide {i+1} Illustration">
        </div>"""
            
        slide_html = f"""
    <section class="slide" id="slide-{i+1}">
      <div class="slide-layout">
        <div class="slide-text">
          <h2>{title}</h2>
          <ul>
{bullets_html}
          </ul>
        </div>{image_html}
      </div>
    </section>"""
        slides_html_list.append(slide_html)
        
    slides_html = "\n".join(slides_html_list)
    total_slides = len(expanded_slides[:15])
    presentation_title = expanded_slides[0]["title"] if expanded_slides else "學習共同體公開課分析"
    
    html_content = f"""<!DOCTYPE html>
<html lang="zh-TW">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>{presentation_title}</title>
  <style>
    :root {{
      --bg-color: {style_config["bg"]};
      --title-color: {style_config["title"]};
      --text-color: {style_config["text"]};
    }}
    body {{
      margin: 0;
      padding: 0;
      background-color: var(--bg-color);
      font-family: 'Microsoft JhengHei', -apple-system, sans-serif;
      color: var(--text-color);
      display: flex;
      flex-direction: column;
      align-items: center;
      justify-content: center;
      min-height: 100vh;
      overflow: hidden;
    }}
    .slides-container {{
      position: relative;
      width: 90vw;
      max-width: 1200px;
      height: 70vh;
      min-height: 500px;
      background: rgba(255, 255, 255, 0.7);
      backdrop-filter: blur(10px);
      border-radius: 16px;
      box-shadow: 0 10px 30px rgba(0, 0, 0, 0.05);
      border: 1px solid rgba(255, 255, 255, 0.3);
      display: flex;
      flex-direction: column;
      overflow: hidden;
    }}
    .slide {{
      display: none;
      flex: 1;
      padding: 40px;
      box-sizing: border-box;
    }}
    .slide.active {{
      display: flex;
      animation: fadeIn 0.5s ease-in-out forwards;
    }}
    @keyframes fadeIn {{
      from {{ opacity: 0; transform: translateY(10px); }}
      to {{ opacity: 1; transform: translateY(0); }}
    }}
    .slide-layout {{
      display: flex;
      width: 100%;
      height: 100%;
      gap: 40px;
      align-items: center;
    }}
    .slide-text {{
      flex: 1.2;
      display: flex;
      flex-direction: column;
      justify-content: center;
    }}
    .slide-text h2 {{
      color: var(--title-color);
      margin-top: 0;
      margin-bottom: 24px;
      font-size: 32px;
      line-height: 1.3;
      border-left: 5px solid var(--title-color);
      padding-left: 15px;
    }}
    .slide-text ul {{
      margin: 0;
      padding-left: 20px;
    }}
    .slide-text li {{
      font-size: 20px;
      line-height: 1.8;
      margin-bottom: 12px;
      list-style-type: square;
    }}
    .slide-image {{
      flex: 0.8;
      display: flex;
      justify-content: center;
      align-items: center;
      height: 100%;
    }}
    .slide-image img {{
      max-width: 100%;
      max-height: 90%;
      border-radius: 12px;
      box-shadow: 0 5px 20px rgba(0,0,0,0.08);
      object-fit: cover;
      aspect-ratio: 1;
    }}
    .controls {{
      display: flex;
      align-items: center;
      justify-content: space-between;
      padding: 20px 40px;
      background: rgba(255, 255, 255, 0.4);
      border-top: 1px solid rgba(255, 255, 255, 0.2);
    }}
    .btn {{
      background: var(--title-color);
      color: var(--bg-color);
      border: none;
      padding: 10px 20px;
      font-size: 16px;
      font-weight: bold;
      border-radius: 8px;
      cursor: pointer;
      transition: opacity 0.2s, transform 0.1s;
    }}
    .btn:hover {{
      opacity: 0.9;
    }}
    .btn:active {{
      transform: scale(0.98);
    }}
    .btn:disabled {{
      background: #ccc;
      cursor: not-allowed;
    }}
    .progress-bar-container {{
      flex: 1;
      height: 8px;
      background: rgba(0, 0, 0, 0.05);
      border-radius: 4px;
      margin: 0 40px;
      overflow: hidden;
    }}
    .progress-bar {{
      height: 100%;
      background: var(--title-color);
      width: 0%;
      transition: width 0.3s ease;
    }}
    .slide-counter {{
      font-size: 16px;
      font-weight: bold;
      color: var(--text-color);
      min-width: 60px;
      text-align: right;
    }}
    @media (max-width: 768px) {{
      .slide-layout {{
        flex-direction: column;
        overflow-y: auto;
      }}
      .slide-image {{
        height: auto;
        margin-top: 20px;
      }}
      .slide-image img {{
        max-height: 200px;
      }}
    }}
  </style>
</head>
<body>

  <div class="slides-container">
{slides_html}
    
    <!-- Navigation Controls -->
    <div class="controls">
      <button class="btn" id="prev-btn" onclick="changeSlide(-1)">上一頁</button>
      <div class="progress-bar-container">
        <div class="progress-bar" id="progress-bar"></div>
      </div>
      <button class="btn" id="next-btn" onclick="changeSlide(1)">下一頁</button>
      <span class="slide-counter" id="slide-counter">1 / {total_slides}</span>
    </div>
  </div>

  <script>
    let currentSlideIndex = 0;
    const slides = document.querySelectorAll('.slide');
    const totalSlides = slides.length;
    const prevBtn = document.getElementById('prev-btn');
    const nextBtn = document.getElementById('next-btn');
    const progressBar = document.getElementById('progress-bar');
    const slideCounter = document.getElementById('slide-counter');

    function showSlide(index) {{
      slides.forEach((slide, i) => {{
        slide.classList.toggle('active', i === index);
      }});
      
      currentSlideIndex = index;
      
      // Update buttons
      prevBtn.disabled = index === 0;
      nextBtn.disabled = index === totalSlides - 1;
      
      // Update progress bar
      const progressPercent = ((index + 1) / totalSlides) * 100;
      progressBar.style.width = `${{progressPercent}}%`;
      
      // Update counter
      slideCounter.innerText = `${{index + 1}} / ${{totalSlides}}`;
    }}

    function changeSlide(direction) {{
      let targetIndex = currentSlideIndex + direction;
      if (targetIndex >= 0 && targetIndex < totalSlides) {{
        showSlide(targetIndex);
      }}
    }}

    // Keyboard navigation
    document.addEventListener('keydown', (e) => {{
      if (e.key === 'ArrowRight' || e.key === ' ') {{
        changeSlide(1);
      }} else if (e.key === 'ArrowLeft') {{
        changeSlide(-1);
      }}
    }});

    // Initialize first slide
    showSlide(0);
  </script>
</body>
</html>
"""
    output_html.parent.mkdir(parents=True, exist_ok=True)
    output_html.write_text(html_content, encoding="utf-8")
    print(f"[OK] Created HTML slides: {output_html}")
    return 0


def generate_image(text_path: Path, output_png: Path, bg_color: str) -> int:
    """Generates a FB/IG concept image using Pillow."""
    print(f"Generating concept image...")
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("Error: 'Pillow' package is not installed. Please run with 'uv run --with pillow'.", file=sys.stderr)
        return 1

    if not text_path.exists():
        print(f"Error: Text file {text_path} does not exist.", file=sys.stderr)
        return 1

    try:
        lines = [line.strip() for line in text_path.read_text(encoding="utf-8").splitlines() if line.strip()][:8]
    except Exception as e:
        print(f"Error reading text: {e}", file=sys.stderr)
        return 1

    # Image setup: square for IG/FB
    width, height = 1080, 1080
    bg_colors = {
        "pastel": (245, 245, 240),
        "blue": (230, 240, 250),
        "modern": (250, 250, 250),
        "dark": (30, 30, 35),
        "learning": (238, 242, 235)
    }
    text_colors = {
        "pastel": (80, 50, 40),
        "blue": (10, 40, 80),
        "modern": (30, 30, 30),
        "dark": (240, 240, 245),
        "learning": (35, 65, 45)
    }
    
    bg = bg_colors.get(bg_color, bg_colors["modern"])
    text_color = text_colors.get(bg_color, text_colors["modern"])
    
    img = Image.new("RGB", (width, height), color=bg)
    draw = ImageDraw.Draw(img)
    
    # Font setup
    fonts_to_try = [
        "msjh.ttc",            # Microsoft JhengHei on Windows
        "Microsoft JhengHei.ttf",
        "msyh.ttc",            # Microsoft YaHei
        "arial.ttf"
    ]
    font_title = None
    font_body = None
    
    for f_name in fonts_to_try:
        try:
            font_title = ImageFont.truetype(f_name, 48)
            font_body = ImageFont.truetype(f_name, 28)
            break
        except IOError:
            continue
            
    if not font_title:
        # Fallback to default
        font_title = ImageFont.load_default()
        font_body = ImageFont.load_default()

    # Draw card outline
    draw.rectangle([60, 60, width - 60, height - 60], outline=text_color, width=3)
    
    # Draw Title
    title = "學習共同體 • 課堂核心概念"
    draw.text((100, 120), title, fill=text_color, font=font_title)
    draw.line([100, 190, width - 100, 190], fill=text_color, width=2)
    
    # Draw Bullet Points
    y_offset = 240
    for line in lines:
        if len(line) > 35:
            line = line[:35] + "..."
        draw.text((120, y_offset), "★ " + line, fill=text_color, font=font_body)
        y_offset += 75
        
    # Draw Footer
    footer = "光復國小 數學學科公開課 課例研究成果分享"
    draw.text((100, height - 120), footer, fill=text_color, font=font_body)

    output_png.parent.mkdir(parents=True, exist_ok=True)
    img.save(str(output_png))
    print(f"[OK] Generated concept image: {output_png}")
    return 0

def main() -> int:
    parser = argparse.ArgumentParser(description="Classroom Video Analyzer Helper CLI")
    subparsers = parser.add_subparsers(dest="command", required=True)
    
    # fetch-audio
    p_fetch = subparsers.add_parser("fetch-audio", help="Fetch audio track from URL as mp3")
    p_fetch.add_argument("url", type=str, help="YouTube or video URL")
    p_fetch.add_argument("--output", type=Path, required=True, help="Destination mp3 file path")
    
    # transcribe
    p_trans = subparsers.add_parser("transcribe", help="Transcribe audio to SRT/TXT using Whisper")
    p_trans.add_argument("audio", type=Path, help="Source audio file path")
    p_trans.add_argument("--output-srt", type=Path, required=True, help="Destination SRT subtitle path")
    p_trans.add_argument("--output-txt", type=Path, required=True, help="Destination TXT transcript path")
    p_trans.add_argument("--model", type=str, default="medium", help="Whisper model size")
    p_trans.add_argument("--device", type=str, default="cpu", help="Computation device (cpu or cuda)")
    
    # generate-slides
    p_slides = subparsers.add_parser("generate-slides", help="Generate PPTX slides from analysis text")
    p_slides.add_argument("--analysis", type=Path, required=True, help="Path to analysis markdown/text file")
    p_slides.add_argument("--output", type=Path, required=True, help="Destination PPTX file path")
    p_slides.add_argument("--style", type=str, default="modern", choices=["modern", "pastel", "blue", "learning"], help="Slides color theme")
    
    # generate-image
    p_image = subparsers.add_parser("generate-image", help="Generate social media conceptual image")
    p_image.add_argument("--text", type=Path, required=True, help="Path to text summary file")
    p_image.add_argument("--output", type=Path, required=True, help="Destination PNG file path")
    p_image.add_argument("--style", type=str, default="modern", choices=["modern", "pastel", "blue", "dark", "learning"], help="Image color style")
    
    args = parser.parse_args()
    
    if args.command == "fetch-audio":
        return fetch_audio(args.url, args.output)
    elif args.command == "transcribe":
        return transcribe_audio(args.audio, args.output_srt, args.output_txt, args.model, args.device)
    elif args.command == "generate-slides":
        return generate_slides(args.analysis, args.output, args.style)
    elif args.command == "generate-image":
        return generate_image(args.text, args.output, args.style)
        
    return 0

if __name__ == "__main__":
    sys.exit(main())
